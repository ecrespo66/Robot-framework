from pptx.util import Pt
from iBott.files_and_folders.files import File
import os
from pptx import Presentation
from pptx.slide import Slide


def add_title(self, text, size):
    title = self.shapes.title
    title.text = text
    title.size = Pt(size)


def add_subtitle(self, text, size):
    subtitle = self.placeholders[1]
    subtitle.text = text
    subtitle.size = Pt(size)


def add_text(cls, text, size, shape=(100, 100, 100, 100)):
    txBox = cls.shapes.add_textbox(shape[0], shape[1], shape[2], shape[3])
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.size = Pt(size)


def add_image(cls, image, shape=(100, 100, 100, 100)):
    cls.shapes.add_picture(image, shape[0], shape[1], shape[2], shape[3])


class PowerPoint(File):
    """
    Class to manage powerpoint files
    it heritages from File class
    Arguments:
        file_path (str): path to the file
    Methods:
        add_slide(layout): add a slide to the presentation
        get_slide(index): get a slide from the presentation
        get_slides(self)
    """

    def __init__(self, file_path):
        self.ppt = self.__file(file_path)
        super().__init__(file_path)
        setattr(Slide, 'add_title', add_title)
        setattr(Slide, 'add_subtitle', add_subtitle)
        setattr(Slide, 'add_text', add_text)

    def add_slide(self, layout=6):
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[layout])
        self.ppt.save(self.path)
        return slide

    def get_slide(self, index):
        return self.ppt.slides[index]

    def get_slides(self):
        return self.ppt.slides

    def get_slide_count(self):
        return len(self.ppt.slides)

    def get_slide_layout(self):
        return self.ppt.slides[0].slide_layout

    def save(self):
        self.ppt.save(self.path)

    @staticmethod
    def __file(path):
        if os.path.isfile(path):
            return Presentation(path)
        else:
            ppt = Presentation()
            ppt.save(path)
        return ppt

    @classmethod
    def add_title(cls, text):
        cls.shapes.title.text = text

    @classmethod
    def add_subtitle(cls, text):
        cls.placeholders[1].text = text




