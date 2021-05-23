from iBott.files_activities import File
import os
from pptx import Presentation
from pptx.slide import Slide


def add_title(self, text):
    self.shapes.title.text = text


def add_subtitle(self, text):
    self.placeholders[1].text = text


class PowerPoint(File):
    def __init__(self, path):
        self.ppt = self.__file(path)
        super().__init__(path)
        setattr(Slide, 'add_title', add_title)

    @staticmethod
    def __file(path):
        if os.path.isfile(path):
            return Presentation(path)
        else:
            ppt = Presentation()
            ppt.save(path)
        return ppt

    def add_slide(self, layout):
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[layout])
        self.ppt.save(self.path)
        return slide

    def save(self):
        self.ppt.save(self.path)


ppt = PowerPoint('/Users/enriquecrespodebenito/Desktop/pdfDeclaracion.pptx')

slide = ppt.add_slide(2)
slide.add_title('POOLS')
ppt.save()
